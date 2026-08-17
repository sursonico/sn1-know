"""
kb/ingest.py — Ingestion pipeline for the SN1 Knowledge Base.

Usage:
    python -m kb.ingest                    # ingest all new files in DOCS_DIR
    python -m kb.ingest path/to/file.pdf   # ingest a single file
"""

import asyncio
import hashlib
import io
import logging
import re
import shutil
import sys
from dataclasses import dataclass, field
from pathlib import Path
from typing import Optional

import pdfplumber
from pptx import Presentation

from config import (
    DOCS_DIR, DB_PATH, INGEST_BATCH_SIZE, DEAL_EXTRACT_CONCURRENCY,
    MAX_CHARS_PER_DOC, MAX_CHARS_PER_CHUNK, OCR_CHAR_THRESHOLD,
    VISION_CHAR_THRESHOLD, VISION_IMAGE_COUNT_THRESHOLD, VISION_HYBRID_MAX_CHARS,
    VALIDATION_MIN_FILE_BYTES, VALIDATION_MIN_CHARS_PER_KB,
    VALIDATION_MIN_CHUNKS_FOR_ENTITY_CHECK,
)
from kb import db
from kb.files import store_path, clear_lookup_cache
from kb.llm import (
    classify_document_async, resolve_entities_async, extract_deals_async,
    vision_sdk_available, describe_page_images_async,
)

log = logging.getLogger("sn1.ingest")


# ── Dataclass for a single page/slide/sheet chunk ────────────────────────────

@dataclass
class Chunk:
    chunk_num: int
    chunk_type: str   # 'page' | 'slide' | 'sheet' | 'body'
    text: str
    image_count: int = 0  # embedded images/pictures on this page or slide


@dataclass
class ExtractionResult:
    file_type: str
    chunks: list[Chunk] = field(default_factory=list)
    ocr_used: bool = False
    error: str = ""


# ── Text extraction ───────────────────────────────────────────────────────────

_COLUMN_GAP_PT = 18.0  # x0 gap that marks a column boundary, empirically tuned
                       # against the PROPERTIES OVERVIEW deck's 5-column layout


def _detect_column_edges(words: list[dict], page_width: float) -> Optional[list[float]]:
    """Return x-axis column boundaries, or None if the page reads as one column."""
    if not words:
        return None
    xs_sorted = sorted(w["x0"] for w in words)
    gaps = [(a, b) for a, b in zip(xs_sorted, xs_sorted[1:]) if b - a > _COLUMN_GAP_PT]
    if not gaps:
        return None
    edges = [0.0] + [(a + b) / 2 for a, b in gaps] + [page_width]
    # A wide final gap can push its midpoint past the page edge; page.crop()
    # raises on a bbox outside page bounds, so clamp it back in.
    edges[-1] = min(edges[-1], page_width)
    return edges


def _extract_pdf_page_text(page) -> str:
    """
    Extract page text in left-to-right column order instead of raw word
    position. Side-by-side sections (e.g. "Overview | Audience | Media rights
    deals | Analysis") read out of order under plain extract_text(), which
    walks words top-to-bottom across the whole page width and splices
    unrelated columns mid-sentence. Detects column boundaries as gaps
    >_COLUMN_GAP_PT in the sorted x0 distribution, crops each band, and
    concatenates column-by-column.
    """
    try:
        words = page.extract_words()
    except Exception:
        words = []

    edges = _detect_column_edges(words, page.width)
    if not edges:
        return (page.extract_text() or "").strip()

    bands = [[lo, hi, [w for w in words if lo <= w["x0"] < hi]] for lo, hi in zip(edges, edges[1:])]

    # A sliver band (a stray right-aligned page number, a logo) isn't a real
    # column — merge it into the previous band rather than cropping it alone,
    # which would scramble reading order instead of fixing it.
    merged: list[list] = []
    for lo, hi, band_words in bands:
        if merged and len(band_words) < 3:
            merged[-1][1] = hi
            merged[-1][2].extend(band_words)
        else:
            merged.append([lo, hi, band_words])

    if len(merged) < 2:
        return (page.extract_text() or "").strip()

    column_texts = []
    for lo, hi, band_words in merged:
        if not band_words:
            continue
        bbox = (max(lo, 0.0), 0.0, min(hi, page.width), page.height)
        try:
            text = (page.crop(bbox).extract_text() or "").strip()
        except Exception:
            text = ""
        if text:
            column_texts.append(text)

    return "\n\n".join(column_texts) if column_texts else (page.extract_text() or "").strip()


def _extract_pdf(path: Path) -> ExtractionResult:
    chunks: list[Chunk] = []
    scanned_pages: list[int] = []
    try:
        with pdfplumber.open(path) as pdf:
            for i, page in enumerate(pdf.pages, start=1):
                text = _extract_pdf_page_text(page)
                n_images = len(page.images)
                if len(text) < OCR_CHAR_THRESHOLD:
                    scanned_pages.append(i)
                    text = text or f"[page {i} — no extractable text]"
                chunks.append(Chunk(i, "page", text[:MAX_CHARS_PER_CHUNK], image_count=n_images))
    except Exception as e:
        return ExtractionResult("PDF", error=str(e))

    ocr_used = False
    if len(scanned_pages) > len(chunks) // 2 + 1:
        ocr_used = _try_ocr_fallback(path, chunks, scanned_pages)

    return ExtractionResult("PDF", chunks=chunks, ocr_used=ocr_used)


def _try_ocr_fallback(path: Path, chunks: list[Chunk], scanned_pages: list[int]) -> bool:
    """Replace sparse-text chunks with OCR text. Returns True if OCR was applied."""
    try:
        import pytesseract
        from pdf2image import convert_from_path
        images = convert_from_path(str(path), dpi=200)
        for i in scanned_pages:
            if i - 1 < len(images):
                ocr_text = pytesseract.image_to_string(images[i - 1]).strip()
                if ocr_text and i - 1 < len(chunks):
                    chunks[i - 1] = Chunk(i, "page", ocr_text[:MAX_CHARS_PER_CHUNK])
        return True
    except ImportError:
        return False
    except Exception as e:
        log.warning("OCR failed for %s: %s", path.name, e)
        return False


# ── Vision enrichment for image-heavy pages ───────────────────────────────────

def _render_pdf_page_image(path: Path, page_idx: int) -> Optional[bytes]:
    """
    Render a PDF page as PNG bytes.
    Uses pymupdf (fitz) if installed — full rendering including embedded images.
    Falls back to pdfplumber.to_image() — renders text/vector only; warns if
    the page appears image-only (embedded images but nearly blank render).
    Returns None if rendering is not possible.
    """
    try:
        import fitz  # pymupdf
        doc = fitz.open(str(path))
        if page_idx >= len(doc):
            return None
        page = doc[page_idx]
        pix = page.get_pixmap(matrix=fitz.Matrix(1.5, 1.5))
        return pix.tobytes("png")
    except ImportError:
        pass  # fall through to pdfplumber
    except Exception as e:
        log.debug("pymupdf render failed for %s page %d: %s", path.name, page_idx, e)

    # Fallback: pdfplumber.to_image() — renders vector/text layer only
    try:
        buf = io.BytesIO()
        with pdfplumber.open(path) as pdf:
            if page_idx >= len(pdf.pages):
                return None
            page = pdf.pages[page_idx]
            n_embedded = len(page.images)
            img = page.to_image(resolution=120)
            img.original.save(buf, format="PNG")
            png = buf.getvalue()
            if n_embedded > 0 and len(png) < 15_000:
                log.warning(
                    "%s page %d: %d embedded images but pdfplumber rendered only %d bytes "
                    "(likely image-only page). Install pymupdf for full vision: pip install pymupdf",
                    path.name, page_idx + 1, n_embedded, len(png),
                )
            return png
    except Exception as e:
        log.debug("pdfplumber.to_image failed for %s page %d: %s", path.name, page_idx, e)

    return None


def _extract_pptx_slide_images(path: Path, slide_idx: int) -> tuple:
    """
    Extract embedded image blobs from a PPTX slide (PNG and JPEG/WEBP only; skips WMF/EMF).
    Also captures text from the slide for context.
    Returns (image_blobs, context_text).
    """
    try:
        from pptx import Presentation

        prs = Presentation(path)
        if slide_idx >= len(prs.slides):
            return [], ""

        slide = prs.slides[slide_idx]
        images: list = []
        texts: list = []

        def _collect(shapes):
            sorted_shapes = sorted(
                shapes,
                key=lambda s: (getattr(s, "top", 0), getattr(s, "left", 0)),
            )
            for shape in sorted_shapes:
                if shape.shape_type == 13:  # MSO_SHAPE_TYPE.PICTURE
                    try:
                        blob = shape.image.blob
                        if (
                            blob[:4] == b'\x89PNG'
                            or blob[:2] == b'\xff\xd8'
                            or (blob[:4] == b'RIFF' and len(blob) > 12 and blob[8:12] == b'WEBP')
                        ):
                            images.append(blob)
                    except Exception:
                        pass
                elif shape.shape_type == 6:  # MSO_SHAPE_TYPE.GROUP
                    try:
                        _collect(shape.shapes)
                    except Exception:
                        pass
                elif shape.has_text_frame:
                    try:
                        lines = [
                            " ".join(r.text for r in p.runs).strip()
                            for p in shape.text_frame.paragraphs
                        ]
                        text = " ".join(l for l in lines if l)
                        if text:
                            texts.append(text)
                    except Exception:
                        pass

        _collect(slide.shapes)
        return images, "\n".join(texts)

    except Exception as e:
        log.debug("PPTX image extraction failed for %s slide %d: %s", path.name, slide_idx, e)
        return [], ""


def _should_vision_chunk(chunk: Chunk) -> bool:
    """
    Returns True if a chunk should be sent to Claude vision.
    Two triggers:
      1. Sparse text: very little text extracted (likely image-only page).
      2. Hybrid: enough embedded images alongside modest text (e.g. broadcaster
         logo tables where country names extract as text but logos do not).
    """
    text_len = len(chunk.text.strip())
    if text_len < VISION_CHAR_THRESHOLD:
        return True
    if (chunk.image_count >= VISION_IMAGE_COUNT_THRESHOLD
            and text_len < VISION_HYBRID_MAX_CHARS):
        return True
    return False


async def _vision_enrich_async(path: Path, result: ExtractionResult) -> ExtractionResult:
    """
    For chunks with sparse text OR significant embedded images alongside modest text,
    attempt Claude vision analysis. Augments chunk text with the vision output
    (marked with [Vision-extracted]).

    Requires ANTHROPIC_API_KEY — logs a warning and returns unchanged result if absent.
    Only runs on PDF and PPTX files.
    """
    if path.suffix.lower() not in (".pdf", ".pptx"):
        return result

    targets = [
        i for i, c in enumerate(result.chunks)
        if _should_vision_chunk(c)
    ]
    if not targets:
        return result

    if not vision_sdk_available():
        log.warning(
            "%s has %d image-heavy page(s) but ANTHROPIC_API_KEY is not set — "
            "vision analysis requires the SDK. Set ANTHROPIC_API_KEY to enable.",
            path.name, len(targets),
        )
        return result

    n_sparse   = sum(1 for i in targets if len(result.chunks[i].text.strip()) < VISION_CHAR_THRESHOLD)
    n_hybrid   = len(targets) - n_sparse
    log.info(
        "%s: vision analysis on %d page(s) (%d sparse, %d hybrid logo/image)",
        path.name, len(targets), n_sparse, n_hybrid,
    )
    new_chunks = list(result.chunks)
    ext = path.suffix.lower()

    for chunk_idx in targets:
        chunk = result.chunks[chunk_idx]
        page_num = chunk.chunk_num   # 1-based
        page_idx = page_num - 1     # 0-based

        if ext == ".pdf":
            img_bytes = _render_pdf_page_image(path, page_idx)
            if not img_bytes:
                continue
            image_blobs = [img_bytes]
            context = chunk.text.strip() or f"PDF page {page_num}"
        else:  # .pptx
            image_blobs, slide_text = _extract_pptx_slide_images(path, page_idx)
            if not image_blobs:
                continue
            context = slide_text or chunk.text.strip() or f"Slide {page_num}"

        try:
            description = await describe_page_images_async(image_blobs, context_hint=context)
        except Exception as e:
            log.warning("Vision call failed for %s page %d: %s", path.name, page_num, e)
            continue

        if not description:
            continue

        existing = chunk.text.strip()
        if existing and existing != f"[page {page_num} — no extractable text]":
            augmented = f"{existing}\n\n[Vision-extracted]\n{description}"
        else:
            augmented = f"[Vision-extracted]\n{description}"

        new_chunks[chunk_idx] = Chunk(
            chunk.chunk_num,
            chunk.chunk_type,
            augmented[:MAX_CHARS_PER_CHUNK],
        )
        log.info("%s page/slide %d: vision added %d chars", path.name, page_num, len(description))

    return ExtractionResult(result.file_type, new_chunks, result.ocr_used)


def _count_pptx_pictures(shapes) -> int:
    count = 0
    for shape in shapes:
        if shape.shape_type == 13:  # MSO_SHAPE_TYPE.PICTURE
            count += 1
        elif shape.shape_type == 6:  # MSO_SHAPE_TYPE.GROUP
            try:
                count += _count_pptx_pictures(shape.shapes)
            except Exception:
                pass
    return count


def _collect_pptx_shape_text(shapes) -> list[str]:
    """Recursively collect paragraph text from shapes, descending into groups.

    Many decks lay out slide body content (stat callouts, bullet lists) inside
    grouped shapes rather than directly on the slide — a flat top-level scan
    only picks up the slide title and leaves everything else unextracted.

    Shapes are sorted by position (.left, then .top) at each recursion level
    before their text is collected — raw shape/document order does not match
    reading order in multi-column layouts, splicing side-by-side sections
    together the same way an unsorted PDF word list does.
    """
    texts: list[str] = []
    try:
        shapes = sorted(shapes, key=lambda s: (getattr(s, "left", None) or 0, getattr(s, "top", None) or 0))
    except Exception:
        pass
    for shape in shapes:
        if shape.shape_type == 6:  # MSO_SHAPE_TYPE.GROUP
            try:
                texts.extend(_collect_pptx_shape_text(shape.shapes))
            except Exception:
                pass
        elif shape.has_text_frame:
            for para in shape.text_frame.paragraphs:
                line = " ".join(r.text for r in para.runs).strip()
                if line:
                    texts.append(line)
    return texts


def _extract_pptx(path: Path) -> ExtractionResult:
    chunks: list[Chunk] = []
    try:
        prs = Presentation(path)
        for i, slide in enumerate(prs.slides, start=1):
            texts = _collect_pptx_shape_text(slide.shapes)
            n_images = _count_pptx_pictures(slide.shapes)
            chunks.append(Chunk(i, "slide", "\n".join(texts)[:MAX_CHARS_PER_CHUNK], image_count=n_images))
    except Exception as e:
        return ExtractionResult("PowerPoint", error=str(e))
    return ExtractionResult("PowerPoint", chunks=chunks)


def _extract_xlsx(path: Path) -> ExtractionResult:
    chunks: list[Chunk] = []
    try:
        import openpyxl
        wb = openpyxl.load_workbook(path, read_only=True, data_only=True)
        for i, ws in enumerate(wb.worksheets, start=1):
            rows_text: list[str] = []
            for j, row in enumerate(ws.iter_rows(values_only=True)):
                if j >= 200:
                    break
                cells = [str(c) for c in row if c is not None]
                if cells:
                    rows_text.append("\t".join(cells))
            chunks.append(Chunk(i, "sheet", "\n".join(rows_text)[:MAX_CHARS_PER_CHUNK]))
        wb.close()
    except Exception as e:
        return ExtractionResult("Excel", error=str(e))
    return ExtractionResult("Excel", chunks=chunks)


_EXTRACTORS = {
    ".pdf":  _extract_pdf,
    ".pptx": _extract_pptx,
    ".xlsx": _extract_xlsx,
    ".xls":  _extract_xlsx,
}

FILE_TYPE_LABELS = {
    ".pdf": "PDF", ".pptx": "PowerPoint", ".xlsx": "Excel", ".xls": "Excel",
}


def extract(path: Path) -> ExtractionResult:
    ext = path.suffix.lower()
    fn  = _EXTRACTORS.get(ext)
    if fn is None:
        return ExtractionResult(ext.lstrip(".").upper(), error=f"Unsupported type: {ext}")
    return fn(path)


def full_text(result: ExtractionResult) -> str:
    """
    Concatenate chunk text for document-level classification, capped at
    MAX_CHARS_PER_DOC. Each chunk gets a fair per-chunk share of the budget
    before joining, so a many-page/slide document doesn't have its back half
    squeezed out entirely by front matter that alone fills the whole cap —
    every page still contributes something to classification and entity/deal
    extraction.
    """
    texts = [c.text for c in result.chunks if c.text.strip()]
    if not texts:
        return ""
    per_chunk = max(MAX_CHARS_PER_DOC // len(texts), 200)
    return "\n\n".join(t[:per_chunk] for t in texts)[:MAX_CHARS_PER_DOC]


# ── Content hashing ───────────────────────────────────────────────────────────

def content_hash(path: Path) -> str:
    h = hashlib.sha256()
    with path.open("rb") as f:
        for block in iter(lambda: f.read(65_536), b""):
            h.update(block)
    return h.hexdigest()


# ── Duplicate / version detection ────────────────────────────────────────────

_VERSION_RE = re.compile(
    r"[-_](v\d+|final|draft|revised|updated|new|old|backup)[\s._]",
    re.IGNORECASE,
)

def looks_like_version(name: str) -> bool:
    stem = Path(name).stem
    return bool(_VERSION_RE.search(stem)) or bool(re.search(r"\(\d+\)$", stem))


def find_possible_duplicate(source: str, existing_sources: list[str]) -> Optional[str]:
    """Return an existing source name that looks like the same file under a different version."""
    stem = Path(source).stem.lower()
    # Normalise away version tokens
    normalised = re.sub(r"[-_](v\d+|final|draft|revised|updated|new|old|\d{4}[-_]\d{2}[-_]\d{2})", "", stem)
    for existing in existing_sources:
        e_stem = Path(existing).stem.lower()
        e_norm = re.sub(r"[-_](v\d+|final|draft|revised|updated|new|old|\d{4}[-_]\d{2}[-_]\d{2})", "", e_stem)
        if e_norm == normalised and existing != source:
            return existing
    return None


# ── Post-ingest validation (advisory — never blocks ingestion) ───────────────

_BULLET_CURRENCY_RE = re.compile(r"^[ \t]*[•\-][^\n]*[$€£]\s?\d", re.MULTILINE)


def _compute_validation_warning(
    chunk_texts: list[str],
    file_size_bytes: int,
    n_entities: int,
    chunk_deal_counts: Optional[list[tuple[int, str, int]]] = None,
) -> str:
    """
    Flag likely-incomplete ingestion so it can be surfaced for human review.
    Never raises and never affects ingestion — a purely advisory signal.

    Takes plain chunk text (not an ExtractionResult) so the Add & Log manual
    review UI can re-run this exact check after a human edits chunk text or
    adds entity links, using data read straight back from the DB — see
    pages/add_log.py's "Save corrections" flow.

    Three independent checks, joined if more than one fires:
      1. Thin extraction: total extracted chars are low relative to file size.
         Skipped below VALIDATION_MIN_FILE_BYTES so small, legitimately sparse
         files (a one-line note, a short single-page PDF) aren't flagged.
      2. Low entity yield: a document with enough pages/slides to plausibly be
         multi-topic (>= VALIDATION_MIN_CHUNKS_FOR_ENTITY_CHECK) produced far
         fewer distinct entities than sections — the exact pattern that hid the
         SPORTS PROPERTIES OVERVIEW group-shape extraction bug (13 slides, 0
         entities).
      3. Silent per-chunk deal-extraction shortfall: `chunk_deal_counts` is
         (chunk_num, text, n_deals) for every chunk extraction was actually
         attempted on. Counts currency-bearing bullet lines in the text as a
         proxy for how many deals the page plausibly has, and flags any chunk
         where the written count is well below that — not just zero. A page
         losing 18 of 19 deals (Serie A: 19 bullets, 1 written) is exactly as
         damaging as one losing all of them, and a zero-only check misses it
         entirely. Checking per-chunk rather than a document-level total also
         matters on its own: a page that fails out of a 13-page deck still
         leaves the *document* total well above zero, so an aggregate check
         never fires even though that one page's deals are silently missing —
         what happened to Formula 1/NHL/WTA in PROPERTIES OVERVIEW.pdf's first
         re-ingest. Only fires on chunks dense enough (>=4 bullets) to trust
         the signal, and only when the shortfall is severe (written count
         under half the bullet count) — a looser bar than the retry-time
         check in kb.llm._looks_incomplete, since this is the last line of
         defense after retries are already exhausted.
    """
    warnings: list[str] = []
    n_chunks = len(chunk_texts)
    total_chars = sum(len(t) for t in chunk_texts)

    if file_size_bytes >= VALIDATION_MIN_FILE_BYTES:
        chars_per_kb = total_chars / (file_size_bytes / 1024)
        if chars_per_kb < VALIDATION_MIN_CHARS_PER_KB:
            warnings.append(
                f"Thin extraction: {total_chars:,} chars from a "
                f"{file_size_bytes/1024:.0f}KB file ({chars_per_kb:.0f} chars/KB)"
            )

    if n_chunks >= VALIDATION_MIN_CHUNKS_FOR_ENTITY_CHECK:
        expected_min = max(1, n_chunks // 6)
        if n_entities < expected_min:
            warnings.append(
                f"Only {n_entities} entit{'y' if n_entities == 1 else 'ies'} linked "
                f"across {n_chunks} slides/pages"
            )

    if chunk_deal_counts:
        flagged = []
        for chunk_num, text, n_deals in chunk_deal_counts:
            n_bullets = len(_BULLET_CURRENCY_RE.findall(text))
            if n_bullets >= 4 and n_deals < max(1, n_bullets // 2):
                flagged.append(f"{chunk_num} ({n_deals}/{n_bullets})")
        if flagged:
            warnings.append(
                f"Deal count implausibly low vs. currency bullets on page(s)/slide(s) "
                f"{', '.join(flagged)} (written/bullets)"
            )

    if not warnings:
        return ""
    return "⚠ " + " · ".join(warnings) + " — may be incomplete extraction, review before relying on it."


# ── Single-file ingestion ─────────────────────────────────────────────────────

async def ingest_file_async(path: Path, existing_sources: list[str]) -> str:
    """
    Ingest one file asynchronously. Returns a status string.
    Does NOT commit the DB — the caller batches commits.
    """
    source = path.name

    # Skip if content hash already in DB
    h = content_hash(path)
    if db.hash_exists(h):
        return f"SKIP  {source} (unchanged)"

    # Ensure the file lives inside DOCS_DIR so file_path is always managed.
    # If it came from an external path (e.g. CLI with an absolute arg), copy it in.
    if path.resolve().parent != DOCS_DIR.resolve():
        DOCS_DIR.mkdir(parents=True, exist_ok=True)
        dest = DOCS_DIR / path.name
        if not dest.exists():
            shutil.copy2(str(path), str(dest))
        path = dest

    # Extract text chunks
    result = extract(path)
    if result.error:
        entry_id = db.upsert_document(
            source=source, file_type=result.file_type,
            ingest_error=result.error, file_path=store_path(path),
            content_hash=h,
        )
        log.warning("Extraction error for %s: %s", source, result.error)
        return f"ERROR {source}: {result.error}"

    # Vision enrichment: fill in sparse/image-heavy pages with Claude vision
    result = await _vision_enrich_async(path, result)

    txt = full_text(result)
    if not txt.strip():
        entry_id = db.upsert_document(
            source=source, file_type=result.file_type,
            ingest_error="No text extracted",
            file_path=store_path(path), content_hash=h,
        )
        return f"EMPTY {source}"

    # Classify with Haiku (async)
    try:
        meta = await classify_document_async(source, txt)
    except Exception as e:
        meta = {k: "" for k in ["sports_leagues","time_period","doc_type","notes","summary","topics","org_tags","market_tags"]}
        meta["notes"] = f"Classification failed: {e}"
        log.warning("Classification failed for %s: %s", source, e)

    # Duplicate / version flag
    is_dup = 0
    dup_of = find_possible_duplicate(source, existing_sources)
    if dup_of or looks_like_version(source):
        is_dup = 1
        log.info("Possible duplicate/version: %s (matches %s)", source, dup_of or "pattern")

    # Write to DB
    entry_id = db.upsert_document(
        source          = source,
        entry_date      = meta.get("doc_date", meta.get("time_period", "")),
        coverage_period = meta.get("coverage_period", ""),
        file_type       = result.file_type,
        doc_type        = meta.get("doc_type", ""),
        org_tags        = meta.get("org_tags", ""),
        market_tags     = meta.get("market_tags", ""),
        sport_tags      = meta.get("sports_leagues", ""),
        topic_tags      = meta.get("topics", ""),
        summary         = meta.get("summary", ""),
        notes           = meta.get("notes", ""),
        file_path       = store_path(path),
        content_hash    = h,
        is_duplicate    = is_dup,
        ocr_used        = int(result.ocr_used),
        reliability     = meta.get("reliability", "reported"),
    )

    # Store page/slide chunks
    db_chunks = [db.Chunk(c.chunk_num, c.chunk_type, c.text) for c in result.chunks]
    db.store_chunks(entry_id, db_chunks)

    # Update FTS index
    db.index_entry(entry_id)

    # Resolve and link entities (primary/secondary roles from LLM)
    try:
        resolved = await resolve_entities_async(meta, source=source)
        n_entities = 0
        for r in resolved:
            canonical = r.get("canonical", "").strip()
            if canonical:
                eid = db.find_or_create_entity(
                    canonical,
                    r.get("type", "other"),
                    proposed=bool(r.get("is_new", False)),
                )
                db.link_entry_to_entity(entry_id, eid, role=r.get("role", "secondary"))
                n_entities += 1
    except Exception as e:
        log.warning("Entity resolution failed for %s: %s", source, e)
        resolved = []
        n_entities = 0

    # Extract and store structured deals — per chunk, not the truncated
    # whole-document `txt`. `txt` is capped at MAX_CHARS_PER_DOC and shares
    # that budget proportionally across every chunk (see full_text()), so a
    # slide's "Media rights deals" section — which tends to come after
    # Overview/Audience — is often cut off before the model ever sees it.
    # Running per-chunk on the full untruncated text closes that gap, at the
    # cost of one extraction call per page/slide instead of one per document.
    # Concurrency is capped (DEAL_EXTRACT_CONCURRENCY) rather than firing every
    # chunk at once — with e.g. 13 chunks, unbounded concurrency has been
    # observed to overload the claude-CLI fallback badly enough that calls
    # which would otherwise succeed time out purely from the contention.
    n_deals = 0
    chunk_deal_counts: list[tuple[int, str, int]] = []
    try:
        canonical_names = [
            r.get("canonical", "").strip()
            for r in resolved
            if r.get("canonical", "").strip()
        ]
        if canonical_names:
            sem = asyncio.Semaphore(DEAL_EXTRACT_CONCURRENCY)

            async def _extract_chunk_deals(c):
                async with sem:
                    hint = f"{source} — {c.chunk_type} {c.chunk_num}"
                    return await extract_deals_async(c.text, canonical_names, source_hint=hint)

            per_chunk_deals = await asyncio.gather(*[_extract_chunk_deals(c) for c in result.chunks])
            chunk_deal_counts = [
                (c.chunk_num, c.text, len(raw))
                for c, raw in zip(result.chunks, per_chunk_deals)
            ]
            entry_rel = meta.get("reliability", "reported")
            for raw_deals in per_chunk_deals:
                for d in raw_deals:
                    en = (d.get("entity_name") or "").strip()
                    entity_row = db.find_entity_by_name_or_alias(en)
                    if entity_row:
                        confidence  = d.get("confidence", "medium")
                        deal_status = "unverified" if confidence == "low" else "current"
                        db.add_deal(
                            entity_id       = entity_row["id"],
                            territory       = d.get("territory", ""),
                            broadcaster     = d.get("broadcaster", ""),
                            rights_holder   = d.get("rights_holder", ""),
                            value           = d.get("value"),
                            currency        = d.get("currency", ""),
                            value_note      = d.get("value_note", ""),
                            period_start    = d.get("period_start", ""),
                            period_end      = d.get("period_end", ""),
                            platform        = d.get("platform", ""),
                            source_entry_id = entry_id,
                            source_note     = source,
                            status          = deal_status,
                            reliability     = entry_rel,
                        )
                        n_deals += 1
    except Exception as e:
        log.warning("Deal extraction failed for %s: %s", source, e)

    # Post-ingest validation: advisory only, never blocks — see docstring.
    warning = _compute_validation_warning(
        [c.text for c in result.chunks], path.stat().st_size, n_entities, chunk_deal_counts,
    )
    db.set_validation_warning(entry_id, warning)

    flag = " [possible duplicate]" if is_dup else ""
    flag += " [OCR]" if result.ocr_used else ""
    n_vision = sum(1 for c in result.chunks if "[Vision-extracted]" in c.text)
    flag += f" [vision:{n_vision}]" if n_vision else ""
    flag += f" ({n_entities} entities, {n_deals} deals)"
    flag += " [⚠ NEEDS REVIEW]" if warning else ""
    return f"OK    {source}{flag}"


# ── Batch orchestrator ────────────────────────────────────────────────────────

async def ingest_all_async(
    paths: list[Path],
    batch_size: int = INGEST_BATCH_SIZE,
) -> list[str]:
    """Process `paths` in parallel batches of `batch_size`."""
    db.init_db()
    existing_sources = [e["source"] for e in db.get_all_entries()]
    sem = asyncio.Semaphore(batch_size)

    async def run_one(p: Path) -> str:
        async with sem:
            try:
                return await ingest_file_async(p, existing_sources)
            except Exception as e:
                log.error("Unhandled error for %s: %s", p.name, e)
                return f"FAIL  {p.name}: {e}"

    tasks = [run_one(p) for p in paths]
    results = await asyncio.gather(*tasks)
    clear_lookup_cache()   # new files on disk — drop any cached "not found" lookups
    return results


def ingest_directory(docs_dir: Path = DOCS_DIR) -> list[str]:
    """Synchronous entry-point: ingest all supported files in docs_dir."""
    supported = {".pdf", ".pptx", ".xlsx", ".xls"}
    paths = sorted(
        p for p in docs_dir.iterdir()
        if p.suffix.lower() in supported and not p.name.startswith(".")
    )
    if not paths:
        return [f"No supported files found in {docs_dir}"]
    return asyncio.run(ingest_all_async(paths))


# ── CLI ───────────────────────────────────────────────────────────────────────

if __name__ == "__main__":
    logging.basicConfig(level=logging.INFO, format="%(levelname)s %(message)s")

    if len(sys.argv) > 1:
        targets = [Path(a) for a in sys.argv[1:]]
        results = asyncio.run(ingest_all_async(targets))
    else:
        results = ingest_directory()

    for r in results:
        print(r)

    ok    = sum(1 for r in results if r.startswith("OK"))
    skip  = sum(1 for r in results if r.startswith("SKIP"))
    err   = sum(1 for r in results if r.startswith(("ERROR", "FAIL", "EMPTY")))
    print(f"\n{ok} ingested  {skip} skipped  {err} errors  (total {len(results)})")
